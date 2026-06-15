"""Generate the MMDP project presentation (ATLC_MMDP_presentation.pptx).

Static content (problem, MMDP formulation, algorithms, test scenarios, methodology,
conclusions) is authored here. Results slides are data-driven from the canonical
multi-seed results: the deck uses ``outputs/aggregated/`` (mean +/- std summary table
and aggregated plots with error bands), produced by aggregate_seeds.py over the 5-seed
SLURM sweep. If that directory is missing the script errors out rather than building a
single-run draft -- regenerate the aggregated results first.

Re-run this script after the cluster sweep to regenerate the final deck.

Usage:  python make_presentation.py
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

AGG_DIR = os.path.join("outputs", "aggregated")
OUT_DIR = "outputs"
OUTFILE = "ATLC_MMDP_presentation.pptx"

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREY = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


# --------------------------------------------------------------------------- helpers
def _textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def add_title_slide(title, subtitle, footer):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, Inches(0), Inches(2.4), Inches(13.333), Inches(1.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = _textbox(s, 0.8, 2.55, 11.7, 1.6)
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2 = tf.add_paragraph(); p2.text = subtitle
    p2.font.size = Pt(20); p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    ftf = _textbox(s, 0.8, 6.7, 11.7, 0.5)
    fp = ftf.paragraphs[0]; fp.text = footer
    fp.font.size = Pt(12); fp.font.color.rgb = GREY
    return s


def add_section_header(title):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, Inches(0), Inches(3.1), Inches(13.333), Inches(1.3))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tf = _textbox(s, 0.8, 3.25, 11.7, 1.0)
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return s


def _add_heading(slide, title):
    tf = _textbox(slide, 0.6, 0.35, 12.1, 0.9)
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = NAVY


def add_bullets(title, bullets, body_top=1.35, body_width=12.1, font=18):
    """bullets: list of (text, level) or str."""
    s = prs.slides.add_slide(BLANK)
    _add_heading(s, title)
    tf = _textbox(s, 0.7, body_top, body_width, 7.5 - body_top - 0.3)
    first = True
    for item in bullets:
        text, level = (item, 0) if isinstance(item, str) else item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level
        p.font.size = Pt(font - 2 * level)
        p.font.color.rgb = NAVY if level == 0 else GREY
        if level == 0:
            p.font.bold = True
    return s


def add_image_slide(title, image_path, caption=None, max_w=11.5, max_h=5.4):
    s = prs.slides.add_slide(BLANK)
    _add_heading(s, title)
    if image_path and os.path.exists(image_path):
        from PIL import Image  # pillow ships with python-pptx's deps in practice
        try:
            iw, ih = Image.open(image_path).size
            ratio = min(max_w / (iw / 96.0), max_h / (ih / 96.0))
            w = (iw / 96.0) * ratio
            left = (13.333 - w) / 2
            s.shapes.add_picture(image_path, Inches(left), Inches(1.4), width=Inches(w))
        except Exception:
            s.shapes.add_picture(image_path, Inches(0.9), Inches(1.4), width=Inches(max_w))
    else:
        tf = _textbox(s, 1.0, 3.0, 11.3, 1.0)
        p = tf.paragraphs[0]
        p.text = f"[ plot will appear here after the multi-seed sweep: {os.path.basename(image_path or '')} ]"
        p.font.size = Pt(16); p.font.italic = True; p.font.color.rgb = GREY
    if caption:
        ctf = _textbox(s, 0.7, 6.85, 12.0, 0.5)
        cp = ctf.paragraphs[0]; cp.text = caption
        cp.font.size = Pt(12); cp.font.italic = True; cp.font.color.rgb = GREY
    return s


def add_table_slide(title, header, rows, caption=None, col_widths=None, font=12):
    s = prs.slides.add_slide(BLANK)
    _add_heading(s, title)
    n_rows = len(rows) + 1
    n_cols = len(header)
    gt = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(1.4),
                            Inches(12.1), Inches(0.4 * n_rows)).table
    if col_widths:
        for i, w in enumerate(col_widths):
            gt.columns[i].width = Inches(w)
    for j, h in enumerate(header):
        c = gt.cell(0, j); c.text = str(h)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        para = c.text_frame.paragraphs[0]
        para.font.size = Pt(font); para.font.bold = True
        para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.text = str(val)
            para = c.text_frame.paragraphs[0]
            para.font.size = Pt(font)
            if j > 0:
                para.alignment = PP_ALIGN.RIGHT
    if caption:
        ctf = _textbox(s, 0.7, 6.95, 12.0, 0.4)
        cp = ctf.paragraphs[0]; cp.text = caption
        cp.font.size = Pt(11); cp.font.italic = True; cp.font.color.rgb = GREY
    return s


def parse_summary_md(path):
    """Parse outputs/aggregated/summary.md table into rows of
    (scenario, algo, final, late-window, peak, n seeds)."""
    rows = []
    if not os.path.exists(path):
        return rows
    with open(path, encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line.startswith("|") or "---" in line or "Scenario" in line:
                continue
            cells = [c.strip() for c in line.strip("|").split("|")]
            if len(cells) >= 6:
                rows.append(cells[:6])
    return rows


# --------------------------------------------------------------------------- build deck
add_title_slide(
    "Adaptive Traffic-Signal Control as a Multi-agent MDP",
    "A cooperative MMDP traffic domain and a head-to-head study of value-factorisation MARL",
    "SUMO microsimulation + PettingZoo  ·  IQL · Hysteretic-Q · VDN · QMIX (CTDE)",
)

add_bullets("Motivation & contributions", [
    ("Urban congestion is a coordination problem: fixed-time signals can't adapt to shifting demand,", 0),
    ("and neighbouring junctions must act jointly to move platoons through a corridor.", 1),
    ("We cast adaptive traffic-signal control as a cooperative Multi-agent MDP (MMDP):", 0),
    ("one agent per junction picks a green phase; all agents share a single team reward.", 1),
    ("Contributions", 0),
    ("A configurable MMDP traffic domain: 4 purpose-built stress scenarios + 2 real RESCO benchmarks.", 1),
    ("A head-to-head evaluation of IQL, Hysteretic-Q, VDN and QMIX (CTDE) on it.", 1),
    ("A reproducible 5-seed pipeline reporting mean +/- std, with an open SLURM harness.", 1),
])

# ---- MMDP formulation
add_section_header("Modelling the domain as an MMDP")
add_bullets("MMDP formulation", [
    ("Agents  N", 0),
    ("One controller per signalised junction (1x4: A0,B0,C0,D0; cologne3: 3; grid4x4: 16).", 1),
    ("State  S  (near-full observability -> MMDP, not Dec-POMDP)", 0),
    ("Each agent sees local E-W & N-S queues PLUS the rest-of-network E-W & N-S queues,", 1),
    ("discretised into 5 bins per dimension (5^4 = 625 joint-aware states).", 1),
    ("Actions  A", 0),
    ("Joint action = product of each junction's green-phase choice (1x4: 2 phases; RESCO: up to 8).", 1),
    ("Transitions  T", 0),
    ("SUMO microsimulation with stochastic (Poisson) arrivals and platoon injections.", 1),
    ("Reward  R  (cooperative / shared)", 0),
    ("Single team reward broadcast to all agents: R = - sum of accumulated waiting time over all lanes.", 1),
])
add_bullets("Why this is an MMDP (and why these algorithms fit)", [
    ("MMDP = cooperative multi-agent MDP: a shared/observable joint state and a single team reward.", 0),
    ("The rest-of-network observation pushes each agent toward the global state -> approximates full observability.", 1),
    ("The synchronized global reward makes all agents optimise one cooperative objective.", 1),
    ("This is exactly the setting cooperative value-factorisation methods target:", 0),
    ("VDN and QMIX (CTDE) learn a joint value and factor it for decentralised execution.", 1),
    ("IQL and Hysteretic-Q are decentralised baselines that ignore / partially correct for non-stationarity.", 1),
])

# ---- Algorithms
add_section_header("The algorithms we evaluate")
add_bullets("Independent & Hysteretic Q-Learning (decentralised baselines)", [
    ("Independent Q-Learning (IQL)", 0),
    ("Every agent runs its own tabular Q-learning, treating other agents as part of the environment.", 1),
    ("Simple, but the environment looks non-stationary from each agent's view.", 1),
    ("Hysteretic Q-Learning", 0),
    ("Optimistic IQL: large LR (alpha) for positive TD errors, small LR (beta) for negative ones.", 1),
    ("Reduces mis-penalising an agent for teammates' exploration -> more stable cooperation.", 1),
])
add_bullets("VDN & QMIX (centralised training, decentralised execution)", [
    ("Value Decomposition Networks (VDN)", 0),
    ("Q_tot = sum_i Q_i(s_i, a_i); one shared TD error trains all agents; each acts greedily on its own Q_i.", 1),
    ("Captures the team reward while keeping decentralised execution.", 1),
    ("QMIX", 0),
    ("Q_tot = monotonic mixing network of the Q_i, with weights produced by a hypernetwork on the global state.", 1),
    ("Monotonicity keeps argmax_a Q_tot factorisable; strictly more expressive than VDN's sum.", 1),
    ("Our QMIX is a deep CTDE method (PyTorch agent nets + mixing net) for the 1x4 domain.", 1),
])

# ---- Scenarios (the focus)
add_section_header("Test scenarios: real-world traffic, MMDP-worthy stress tests")
add_table_slide(
    "Custom 1x4 corridor scenarios (we designed these)",
    ["Scenario", "Traffic design", "Real-world analog", "MMDP property stressed"],
    [
        ["baseline", "30-veh platoons / 150 s; cross p=0.1",
         "Arterial green wave, light cross traffic", "Basic spatial coordination"],
        ["dense_wave", "50-veh platoons / 100 s; cross p=0.05",
         "Heavy commuter arterial at peak", "Sequential hand-off (green wave)"],
        ["cross_surge", "20-veh platoons / 180 s; cross p=0.4",
         "Event / shopping-district side streets", "Collective sacrifice of arterial green"],
        ["split_rush", "0-1800s E-W heavy -> 1800-3600s N-S heavy",
         "AM inbound -> PM outbound reversal", "Non-stationary strategy switch"],
    ],
    caption="All four share one 1x4 network; only the demand profile changes, isolating a distinct coordination challenge.",
    col_widths=[1.7, 3.6, 3.4, 3.4], font=12,
)
add_table_slide(
    "RESCO benchmark scenarios (standard, real-world)",
    ["Scenario", "Network", "Real-world analog", "MMDP property stressed"],
    [
        ["cologne3", "Real 3-junction Cologne arterial (OSM), AM rush",
         "Actual city corridor geometry", "Transfer to real, heterogeneous junctions"],
        ["grid4x4", "Synthetic 16-junction grid, right-on-red",
         "Dense downtown grid (Manhattan-like)", "Scalability to many agents"],
    ],
    caption="RESCO = Reinforcement-learning Benchmarks for Signal Control. Junctions have 3-8 green phases (vs 2 in the 1x4).",
    col_widths=[1.6, 4.4, 3.3, 3.0], font=12,
)
add_bullets("Why these scenarios are MMDP-testing-worthy", [
    ("Each scenario isolates a different property of the cooperative joint MDP:", 0),
    ("Spatial credit assignment (dense_wave) - agents must hand platoons downstream in sequence.", 1),
    ("Competing global trade-off (cross_surge) - the team must sacrifice local arterial throughput.", 1),
    ("Temporal non-stationarity (split_rush) - the optimal joint policy changes mid-episode.", 1),
    ("Realism & heterogeneity (cologne3) - real geometry, unequal junctions and phase counts.", 1),
    ("Scale (grid4x4) - 16 interacting agents stress factorisation and exploration.", 1),
    ("Together they probe coordination, not just single-intersection control.", 0),
])

# ---- Methodology
add_section_header("Experimental methodology & reproducibility")
add_bullets("Setup & reproducibility (the \"k-fold\" analog for RL)", [
    ("5 independent seeds x 6 scenarios, run as a SLURM job array (one task per seed x scenario).", 0),
    ("Per-run seeding of numpy / torch / SUMO; fixed per-run eval seed for clean, comparable curves.", 1),
    ("We report mean +/- std across seeds (learning curves with error bands; bar charts with error bars).", 0),
    ("Multi-seed repetition is the RL analog of k-fold cross-validation: it quantifies variance/repro.", 1),
    ("Baselines:", 0),
    ("1x4: coordinated fixed-time (50/50) controller.", 1),
    ("RESCO: BOTH the SUMO-native fixed-time program AND a round-robin all-phases plan.", 1),
    ("Metric: evaluation return = - total system waiting time (\"delay\" = |return|, lower is better).", 0),
])

# ---- Results (data-driven)
add_section_header("Results")
summary_rows = parse_summary_md(os.path.join(AGG_DIR, "summary.md"))
have_agg = bool(summary_rows)

if have_agg:
    add_table_slide(
        "Performance across seeds: final / late-window / peak (mean +/- std)",
        ["Scenario", "Algorithm", "Final", "Late-window", "Peak", "Seeds"],
        summary_rows,
        caption="From outputs/aggregated/summary.md. Return is negative; closer to 0 is better. "
                "Late-window = mean of last 5 evals (fair plateau metric); Peak = best eval per seed.",
        col_widths=[1.7, 2.2, 2.6, 2.6, 2.6, 0.9], font=10,
    )
    add_image_slide("Cross-algorithm comparison: final epoch (mean +/- std over seeds)",
                    os.path.join(AGG_DIR, "cross_algorithm_bars.png"),
                    caption="Final delay = |return| per scenario/algorithm; error bars = std across seeds.")
    late_bars = os.path.join(AGG_DIR, "cross_algorithm_bars_late.png")
    if os.path.exists(late_bars):
        add_image_slide("Cross-algorithm comparison: converged plateau (last 5 evals)",
                        late_bars,
                        caption="Late-window delay = |return|; the fair head-to-head metric, "
                                "robust to late epsilon-greedy noise in the tabular learners.")
    for scen in ["baseline", "dense_wave", "cross_surge", "split_rush", "cologne3", "grid4x4"]:
        lc = os.path.join(AGG_DIR, f"{scen}_learning.png")
        if os.path.exists(lc):
            add_image_slide(f"{scen}: learning curves (mean +/- std)", lc)
        q = os.path.join(AGG_DIR, f"{scen}_qmix.png")
        if os.path.exists(q):
            add_image_slide(f"{scen}: QMIX (CTDE) learning curve", q)

    # ---- Focused analysis of the RESCO / QMIX-at-scale results
    add_bullets("Finding: QMIX (CTDE) extends to real multi-phase junctions", [
        ("We generalised QMIX from the 1x4 topology to RESCO: agent count, observation dim and "
         "per-agent action count are now read from the env (cologne3: 3 agents, 3-4 phases; "
         "grid4x4: 16 agents, 8 phases).", 0),
        ("cologne3 - value factorisation wins decisively:", 0),
        ("All factorised/independent learners beat both fixed-time baselines (native -44k, round-robin -175k).", 1),
        ("QMIX reaches the best single-run policy of any method (peak ~ -2.9k); 4 of 5 seeds converge "
         "to ~ -3.5k.", 1),
        ("But QMIX shows late-training instability: 1 of 5 seeds diverged to -43k, inflating its mean/variance.", 1),
        ("VDN is the most RELIABLE winner on cologne3 (-5.4k +/- 0.4k) - near-best and low variance.", 1),
    ])
    add_bullets("Finding: at 16 agents, expressiveness does NOT buy scalability", [
        ("grid4x4 (16 agents x 8 phases, single global reward): NO learner beats fixed-time (native -31k).", 0),
        ("Among learners, simple VDN is best by far (late-window -156k); every other learner is ~3x worse.", 0),
        ("QMIX is the most UNSTABLE: worst final epoch (-517k) and the highest variance of any learner; its "
         "late-window plateau (-452k) is only mid-pack, and its peak (-307k) is below the tabular learners' best.", 1),
        ("Interpretation: the bottleneck at scale is sample-efficiency / credit assignment / exploration,", 0),
        ("NOT mixer expressiveness. Under a fixed 500-episode budget, QMIX's heavier deep hypernetwork "
         "mixer (over a 64-dim joint state, 16 agents) is harder to train than VDN's simple sum.", 1),
        ("More expressive != better at scale when samples are the binding constraint - a clean negative result.", 1),
    ])
    add_bullets("The scalability picture (complete algorithm x problem grid)", [
        ("Small scale (1x4 corridor, cologne3): value factorisation wins; QMIX/VDN beat fixed-time and "
         "the independent learners.", 0),
        ("Large scale (grid4x4, 16 agents): cooperative MARL hits a wall - even the most expressive "
         "method fails to beat fixed-time, and expressiveness does not rescue it.", 0),
        ("Severity scales with agent count and phase count (2 -> ok, 3-4 -> RL wins, 8x16 -> RL loses),", 1),
        ("pointing to coordination/credit-assignment at scale as the open problem, not the algorithm class.", 1),
    ])
else:
    raise SystemExit(
        "outputs/aggregated/summary.md not found or empty. Run the multi-seed sweep "
        "(cluster/submit.sh) then aggregate_seeds.py before building the deck; the deck "
        "is now always generated from the canonical aggregated multi-seed results."
    )

# ---- Findings / conclusions
add_section_header("Takeaways & outlook")
add_bullets("Takeaways", [
    ("On the 1x4 MMDP, value factorisation wins: QMIX (CTDE) beats coordinated fixed-time on every scenario.", 0),
    ("VDN/QMIX (CTDE) consistently outperform independent learners (IQL/Hysteretic) under a shared reward.", 1),
    ("Generalised to real RESCO junctions: on cologne3 (3 agents) factorisation beats fixed-time by ~10x; "
     "QMIX hits the best peak policy, VDN is the most reliable.", 0),
    ("Scalability wall: at grid4x4 (16 agents) NO learner beats fixed-time; simple VDN is best by far while "
     "QMIX's heavier mixer is the most unstable (worst final epoch, highest variance) - expressiveness does "
     "not buy scale under a fixed sample budget.", 0),
    ("Variance matters: we report final / late-window / peak (mean +/- std across 5 seeds), not single runs; "
     "deep learners can find a great policy yet fail to hold it (1 cologne3 QMIX seed diverged).", 0),
    ("Modelling fidelity is decisive: getting the action space and observation right was the difference "
     "between no learning and learning on the real RESCO benchmarks.", 0),
])
add_bullets("Limitations & future work", [
    ("grid4x4 is unsolved by all methods: the open problem is coordination / credit assignment at 16 agents, "
     "not the algorithm class - try a larger sample budget, reward shaping, or per-agent rewards.", 0),
    ("Stabilise deep CTDE: early-stopping / Double-Q / target smoothing to prevent the late-training "
     "divergence seen on cologne3 QMIX.", 0),
    ("Hyperparameter tuning and longer training for both tabular and deep learners on RESCO.", 0),
    ("Add per-agent metrics (throughput, average delay) alongside the global reward.", 0),
    ("Statistical significance tests across seeds for the headline comparisons.", 0),
])
add_title_slide("Thank you", "Questions?",
                "Reproduce: bash cluster/setup_env.sh  ->  bash cluster/submit.sh  ->  python make_presentation.py")

prs.save(OUTFILE)
n = len(prs.slides._sldIdLst)
print(f"Wrote {OUTFILE} ({n} slides). Aggregated results used: {have_agg}")
